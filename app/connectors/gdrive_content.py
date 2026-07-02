import json
import logging
import zipfile
from html.parser import HTMLParser
from io import BytesIO

logger = logging.getLogger("spoon")

# DOCX/XLSX/PPTX are ZIP containers full of XML. A malicious file can be a
# "zip bomb" (tiny compressed size, huge decompressed size) that exhausts
# memory/CPU when parsed by python-docx/openpyxl/python-pptx. Bail out before
# handing the bytes to those parsers if the archive looks abusive.
_MAX_ZIP_UNCOMPRESSED_BYTES = 200_000_000  # 200 MB decompressed, generous for real documents
_MAX_ZIP_COMPRESSION_RATIO = 100  # decompressed / compressed size
_MAX_PDF_PAGES = 500


def _is_suspicious_zip(data: bytes) -> bool:
    try:
        with zipfile.ZipFile(BytesIO(data)) as archive:
            total_uncompressed = 0
            for info in archive.infolist():
                total_uncompressed += info.file_size
                if total_uncompressed > _MAX_ZIP_UNCOMPRESSED_BYTES:
                    return True
                if info.compress_size and (
                    info.file_size / max(info.compress_size, 1) > _MAX_ZIP_COMPRESSION_RATIO
                ):
                    return True
        return False
    except zipfile.BadZipFile:
        # Not a valid zip at all; let the real parser produce the actual error.
        return False

TEXT_MIME_PREFIXES = ("text/",)
TEXT_MIME_TYPES = {
    "application/json",
    "application/xml",
    "application/javascript",
    "application/x-yaml",
    "application/csv",
    "application/rtf",
}

GOOGLE_EXPORT_FORMATS: dict[str, list[str]] = {
    "application/vnd.google-apps.document": [
        "text/plain",
        "text/html",
        "application/pdf",
    ],
    "application/vnd.google-apps.spreadsheet": [
        "text/csv",
        "application/pdf",
    ],
    "application/vnd.google-apps.presentation": [
        "text/plain",
        "application/pdf",
    ],
    "application/vnd.google-apps.drawing": [
        "application/pdf",
        "image/png",
    ],
    "application/vnd.google-apps.site": [
        "text/html",
    ],
    "application/vnd.google-apps.script": [
        "application/vnd.google-apps.script+json",
    ],
}

SKIP_GOOGLE_MIME_TYPES = {
    "application/vnd.google-apps.folder",
    "application/vnd.google-apps.shortcut",
    "application/vnd.google-apps.form",
    "application/vnd.google-apps.map",
}


class _HTMLTextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self._parts: list[str] = []

    def handle_data(self, data: str) -> None:
        if data.strip():
            self._parts.append(data.strip())

    def text(self) -> str:
        return "\n".join(self._parts)


def is_google_app(mime_type: str) -> bool:
    return mime_type.startswith("application/vnd.google-apps.")


def export_formats_for(mime_type: str) -> list[str]:
    if mime_type in GOOGLE_EXPORT_FORMATS:
        return GOOGLE_EXPORT_FORMATS[mime_type]
    if is_google_app(mime_type):
        return ["application/pdf", "text/plain"]
    return []


def should_skip_mime_type(mime_type: str) -> bool:
    return mime_type in SKIP_GOOGLE_MIME_TYPES


def _decode_text(data: bytes) -> str:
    for encoding in ("utf-8", "utf-16", "latin-1"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _extract_pdf(data: bytes) -> str | None:
    try:
        from pypdf import PdfReader
    except ImportError:
        logger.warning("pypdf not installed; PDF text extraction unavailable")
        return None

    reader = PdfReader(BytesIO(data))
    if len(reader.pages) > _MAX_PDF_PAGES:
        logger.warning(
            "Skipping PDF extraction: %d pages exceeds limit of %d",
            len(reader.pages),
            _MAX_PDF_PAGES,
        )
        return None

    pages = [page.extract_text() or "" for page in reader.pages]
    text = "\n\n".join(page.strip() for page in pages if page.strip())
    return text or None


def _extract_docx(data: bytes) -> str | None:
    if _is_suspicious_zip(data):
        logger.warning("Skipping DOCX extraction: suspicious archive (possible zip bomb)")
        return None

    try:
        from docx import Document as DocxDocument
    except ImportError:
        logger.warning("python-docx not installed; DOCX text extraction unavailable")
        return None

    document = DocxDocument(BytesIO(data))
    paragraphs = [p.text.strip() for p in document.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs) or None


def _extract_xlsx(data: bytes) -> str | None:
    if _is_suspicious_zip(data):
        logger.warning("Skipping XLSX extraction: suspicious archive (possible zip bomb)")
        return None

    try:
        from openpyxl import load_workbook
    except ImportError:
        logger.warning("openpyxl not installed; XLSX text extraction unavailable")
        return None

    workbook = load_workbook(BytesIO(data), read_only=True, data_only=True)
    rows: list[str] = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if cells:
                rows.append(" | ".join(cells))
    return "\n".join(rows) or None


def _extract_pptx(data: bytes) -> str | None:
    if _is_suspicious_zip(data):
        logger.warning("Skipping PPTX extraction: suspicious archive (possible zip bomb)")
        return None

    try:
        from pptx import Presentation
    except ImportError:
        logger.warning("python-pptx not installed; PPTX text extraction unavailable")
        return None

    presentation = Presentation(BytesIO(data))
    slides: list[str] = []
    for slide in presentation.slides:
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if parts:
            slides.append("\n".join(parts))
    return "\n\n".join(slides) or None


def _extract_html(data: bytes) -> str | None:
    parser = _HTMLTextExtractor()
    parser.feed(_decode_text(data))
    text = parser.text()
    return text or None


def extract_text(filename: str, mime_type: str, data: bytes) -> str | None:
    if not data:
        return None

    lowered = mime_type.lower()
    extension = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""

    if lowered.startswith(TEXT_MIME_PREFIXES) or lowered in TEXT_MIME_TYPES:
        text = _decode_text(data).strip()
        return text or None

    if lowered in {"text/html", "application/xhtml+xml"}:
        return _extract_html(data)

    if lowered == "application/pdf" or extension == "pdf":
        return _extract_pdf(data)

    if (
        lowered
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        or extension == "docx"
    ):
        return _extract_docx(data)

    if (
        lowered == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        or extension == "xlsx"
    ):
        return _extract_xlsx(data)

    if (
        lowered
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        or extension == "pptx"
    ):
        return _extract_pptx(data)

    if lowered in {"application/msword", "application/vnd.ms-excel", "application/vnd.ms-powerpoint"}:
        return None

    if lowered.startswith("image/") or lowered.startswith("video/") or lowered.startswith("audio/"):
        return None

    # Last resort for unknown types: try plain-text decode.
    text = _decode_text(data).strip()
    if text and sum(ch.isprintable() or ch in "\n\r\t" for ch in text) / len(text) > 0.85:
        return text
    return None


def supermemory_file_type(mime_type: str, filename: str) -> str | None:
    lowered = mime_type.lower()
    if lowered == "application/vnd.google-apps.document":
        return "google_doc"
    if lowered == "application/vnd.google-apps.spreadsheet":
        return "google_sheet"
    if lowered == "application/vnd.google-apps.presentation":
        return "google_slide"
    if lowered == "application/pdf":
        return "pdf"
    if lowered.startswith("image/"):
        return "image"
    if lowered.startswith("video/"):
        return "video"
    if lowered.startswith("text/") or lowered in TEXT_MIME_TYPES:
        return "text"
    extension = filename.rsplit(".", 1)[-1].lower() if "." in filename else ""
    if extension in {"md", "txt", "csv", "json", "xml", "html"}:
        return "text"
    if extension == "pdf":
        return "pdf"
    return None
